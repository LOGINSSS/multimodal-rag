"""文档入库：提取文本 → 结构化分片 → embedding → 写入 Milvus。

支持：.md / .txt / .docx / .pptx / .pdf（简单文本抽取）/ 图片（OCR + VLM 描述）。
MinerU 的完整版面解析（PDF 表格/公式/图）是后续增强项，见 README 的 TODO。
"""
from __future__ import annotations

import base64
import json
import logging
import re
import shutil
import subprocess
import sys
import tempfile
import uuid
from pathlib import Path
from typing import Dict, List

from langchain_text_splitters import MarkdownHeaderTextSplitter, RecursiveCharacterTextSplitter

from . import config, llm, store

logger = logging.getLogger(__name__)

# ---------- 分片 ----------

_HEADERS = [("#", "h1"), ("##", "h2"), ("###", "h3"), ("####", "h4")]

# 超大块兜底切分（中文友好的分隔符）
_FALLBACK = RecursiveCharacterTextSplitter(
    chunk_size=900,
    chunk_overlap=100,
    separators=["\n\n", "\n", "。", "！", "？", "；", ".", "!", "?", ";", " ", ""],
)


def split_text(text: str) -> List[Dict]:
    """按 Markdown 标题结构化分片；无标题则退化为递归分片。返回 [{text, metadata}]。"""
    md_splitter = MarkdownHeaderTextSplitter(
        headers_to_split_on=_HEADERS, strip_headers=False
    )
    md_docs = md_splitter.split_text(text)

    chunks: List[Dict] = []
    for d in md_docs:
        meta = dict(d.metadata)
        if len(d.page_content) > 1200:
            for sub in _FALLBACK.split_text(d.page_content):
                chunks.append({"text": sub.strip(), "metadata": meta})
        else:
            if d.page_content.strip():
                chunks.append({"text": d.page_content.strip(), "metadata": meta})
    return chunks


# ---------- 文本抽取（按扩展名分派） ----------

def _mineru_exe() -> str:
    """找到 mineru 可执行文件（优先 PATH，退回当前 venv 的 Scripts 目录）。"""
    exe = shutil.which("mineru")
    if exe:
        return exe
    name = "mineru.exe" if sys.platform == "win32" else "mineru"
    candidate = Path(sys.executable).parent / name
    return str(candidate) if candidate.exists() else "mineru"


def _pdf_to_markdown(path: Path) -> str:
    """用 MinerU 把 PDF 转 markdown（版面解析 + OCR + 公式 + 表格）。

    以子进程方式运行，把 PyTorch 模型隔离在独立进程里，不污染 FastAPI 进程。
    产出的 markdown 会继续走 split_text() 的结构化分片。
    """
    workdir = _mkdtemp("mineru_", config.DATA_DIR)
    try:
        cmd = [
            _mineru_exe(),
            "-p", str(path),
            "-o", str(workdir),
            "-b", config.MINERU_BACKEND,
            "-m", config.MINERU_METHOD,
            "-l", config.MINERU_LANG,
        ]
        logger.info("MinerU 解析 %s (backend=%s, method=%s)", path.name, config.MINERU_BACKEND, config.MINERU_METHOD)
        subprocess.run(cmd, check=True, capture_output=True, text=True, timeout=config.MINERU_TIMEOUT)

        stem = path.stem
        matches = list(workdir.rglob(f"{stem}.md"))
        if not matches:
            raise FileNotFoundError(f"MinerU 未产出 markdown：{workdir}")
        return max(matches, key=lambda p: p.stat().st_size).read_text(encoding="utf-8")
    finally:
        shutil.rmtree(workdir, ignore_errors=True)


_HEADING_STYLE_RE = re.compile(r"^(?:heading|标题)\s*(\d+)$", re.IGNORECASE)


def _heading_level(style_name: str) -> int:
    """把 Word 段落样式名映射为 markdown 标题级别（1~6），非标题返回 0。

    兼容英文（Heading 1 / heading1）与中文（标题 1 / 标题1）样式名；
    Title/标题 也按一级标题处理。
    """
    name = (style_name or "").strip()
    m = _HEADING_STYLE_RE.match(name)
    if m:
        return min(int(m.group(1)), 6)
    if name.lower() in {"title", "标题"}:
        return 1
    return 0


def _iter_docx_blocks(parent):
    """按文档顺序产出段落(Paragraph)与表格(Table)块。"""
    from docx.oxml.ns import qn
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    for child in parent.element.body.iterchildren():
        if child.tag == qn("w:p"):
            yield Paragraph(child, parent)
        elif child.tag == qn("w:tbl"):
            yield Table(child, parent)


def _table_to_markdown(table) -> str:
    """python-docx 表格 → markdown 表格字符串。处理合并单元格（同行按 tc 去重）。"""
    rows = []
    for row in table.rows:
        cells, seen = [], set()
        for cell in row.cells:
            if id(cell._tc) in seen:  # 横向合并：同一 tc 重复出现
                continue
            seen.add(id(cell._tc))
            cells.append(" ".join(cell.text.split()))
        rows.append(cells)
    if not rows:
        return ""
    width = max(len(r) for r in rows)
    rows = [r + [""] * (width - len(r)) for r in rows]
    out = ["| " + " | ".join(rows[0]) + " |", "| " + " | ".join(["---"] * width) + " |"]
    out += ["| " + " | ".join(r) + " |" for r in rows[1:]]
    return "\n".join(out)


_TABLE_BAND_MAX_CHARS = 1000


def _split_table_bands(md_table: str, max_chars: int = _TABLE_BAND_MAX_CHARS) -> List[str]:
    """超大 markdown 表格按行分带：每组重复表头+分隔行，避免单 chunk 过大。"""
    lines = md_table.splitlines()
    if len(lines) <= 3:
        return [md_table] if md_table else []
    header, sep, data = lines[0], lines[1], lines[2:]
    bands: List[str] = []
    cur, cur_len = [header, sep], len(header) + len(sep)
    for line in data:
        if cur_len + len(line) > max_chars and len(cur) > 2:
            bands.append("\n".join(cur))
            cur, cur_len = [header, sep], len(header) + len(sep)
        cur.append(line)
        cur_len += len(line)
    if len(cur) > 2:
        bands.append("\n".join(cur))
    return bands


def _extract_docx(path: Path):
    """docx → (markdown正文, 表格chunks)。

    按文档顺序迭代：段落按标题归组（标题转 markdown 走原有结构化切分）；
    每个表格序列化为 markdown 表格并作为原子 chunk（metadata 带当前标题层级
    作上下文），不走字符兜底切分，保证表格不被切碎。
    """
    import docx
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    doc = docx.Document(str(path))
    lines: List[str] = []
    table_chunks: List[Dict] = []
    headers: Dict[int, str] = {}

    for block in _iter_docx_blocks(doc):
        if isinstance(block, Paragraph):
            text = block.text.strip()
            if not text:
                continue
            level = _heading_level(block.style.name if block.style else "")
            if level:
                headers = {k: v for k, v in headers.items() if k < level}
                headers[level] = text
                lines.append(f"{'#' * level} {text}")
            else:
                lines.append(text)
        elif isinstance(block, Table):
            md = _table_to_markdown(block)
            if not md:
                continue
            meta: Dict = {"table": True}
            meta.update({f"h{k}": v for k, v in headers.items()})
            for band in _split_table_bands(md):
                table_chunks.append({"text": band, "metadata": meta})
    return "\n".join(lines), table_chunks


# ---------- PPTX 处理 ----------

_PPTX_IMAGE_MIN_PX = 80  # 小于该尺寸的图片视为装饰，丢弃


def _mkdtemp(prefix: str, dir: Path) -> Path:
    """创建可写的临时目录。

    沙箱环境下 tempfile.mkdtemp 建出的目录带受限 ACL（创建者都写不进），
    这里用普通 mkdir 替代（实测可写）。
    """
    for _ in range(100):
        p = dir / f"{prefix}{uuid.uuid4().hex[:10]}"
        try:
            p.mkdir()
            return p
        except FileExistsError:
            continue
    raise RuntimeError("无法创建临时目录")


def _iter_pptx_shapes(shapes):
    """递归展开 pptx shapes（GROUP 内部也遍历）。"""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_pptx_shapes(shape.shapes)
        else:
            yield shape


def _chart_to_markdown(chart) -> str:
    """python-pptx 图表 → markdown 表格（首列类别，其余列为系列）。"""
    rows = []
    for plot in chart.plots:
        cats = [str(c) for c in plot.categories]
        series = list(plot.series)
        if not series:
            continue
        rows.append(["类别"] + [str(s.name) for s in series])
        for i, cat in enumerate(cats):
            row = [cat]
            for s in series:
                try:
                    row.append(str(s.values[i]))
                except Exception:  # noqa: BLE001 —— 值缺失补空
                    row.append("")
            rows.append(row)
    if not rows:
        return ""
    width = max(len(r) for r in rows)
    rows = [r + [""] * (width - len(r)) for r in rows]
    out = ["| " + " | ".join(rows[0]) + " |", "| " + " | ".join(["---"] * width) + " |"]
    out += ["| " + " | ".join(r) + " |" for r in rows[1:]]
    return "\n".join(out)


def _save_pptx_picture(shape, pptx_path: Path):
    """把 pptx 里的图片导出为临时 PNG 文件，返回 (文件路径, 临时目录)；
    失败（如 WMF/EMF 等 PIL 打不开的格式）返回 (None, None)。"""
    try:
        from PIL import Image

        blob = shape.image.blob
        ext = (shape.image.ext or "png").lower()
        tmp = _mkdtemp("pptx_pic_", config.DATA_DIR)
        try:
            raw = tmp / f"pic_{uuid.uuid4().hex[:8]}.{ext}"
            raw.write_bytes(blob)
            img = Image.open(raw)
            img.load()
            if ext != "png":
                png = tmp / f"{raw.stem}.png"
                img.convert("RGB").save(png)
                return png, tmp
            return raw, tmp
        except Exception:  # noqa: BLE001 —— 单张图片失败不阻断
            shutil.rmtree(tmp, ignore_errors=True)
            return None, None
    except Exception:  # noqa: BLE001
        return None, None


def _picture_chunks(path: Path) -> List[Dict]:
    """单张图片 → [图片描述, 图片OCR] chunks（供入库/pptx 内嵌图片复用）。"""
    ocr_text = _ocr_image(path)
    description = _describe_image(path)
    chunks: List[Dict] = []
    if description:
        chunks.append(
            {
                "text": f"[图片描述] {description}",
                "metadata": json.dumps({"kind": "image_description"}, ensure_ascii=False),
            }
        )
    if ocr_text:
        chunks.append(
            {
                "text": f"[图片OCR文字] {ocr_text}",
                "metadata": json.dumps({"kind": "image_ocr"}, ensure_ascii=False),
            }
        )
    return chunks


def _render_pptx_pages(pptx_path: Path, page_numbers: List[int]) -> Dict[int, str]:
    """用 PowerPoint COM（powershell 子进程）把无文本页渲染成 PNG，再 OCR 出文字。

    尽力而为：渲染或 OCR 失败时该页返回空文本，不影响其他页。
    """
    if not page_numbers:
        return {}
    out_dir = _mkdtemp("pptx_render_", config.DATA_DIR)
    try:
        esc = lambda s: s.replace("'", "''")  # noqa: E731 —— PS 单引号转义
        pages_csv = ",".join(str(n) for n in page_numbers)
        script = (
            "$ErrorActionPreference='Stop'\n"
            "$ppt=New-Object -ComObject PowerPoint.Application\n"
            "try {\n"
            f"  $pres=$ppt.Presentations.Open('{esc(str(pptx_path.resolve()))}')\n"
            "  Start-Sleep -Milliseconds 800\n"
            f"  foreach($n in @({pages_csv})){{ $pres.Slides.Item($n).Export('{esc(str(out_dir))}\\slide_'+$n+'.png','PNG',1280,720) }}\n"
            "  $pres.Close()\n"
            "} finally { $ppt.Quit() }\n"
        )
        subprocess.run(
            ["powershell", "-NoProfile", "-ExecutionPolicy", "Bypass", "-Command", script],
            check=False,
            timeout=180,
        )
    except Exception:  # noqa: BLE001 —— 渲染失败直接放弃兜底
        logger.warning("PPTX 无文本页渲染失败：%s", pptx_path.name)
        return {}
    texts: Dict[int, str] = {}
    try:
        for n in page_numbers:
            png = out_dir / f"slide_{n}.png"
            if png.exists():
                t = _ocr_image(png)
                if t:
                    texts[n] = t
    finally:
        shutil.rmtree(out_dir, ignore_errors=True)
    return texts


def _extract_pptx(path: Path):
    """pptx → (markdown正文, 额外chunks)。

    结构化提取：文本按页组织（`## 第 N 页`）；表格/图表转 markdown 表格 chunk；
    有意义的图片（尺寸够大）走 Qwen-VL 描述 + OCR；无文字的自选图形/连线丢弃；
    整页无任何内容（纯装饰/纯图片且无文本）时用 PowerPoint 渲染 + OCR 兜底。
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(str(path))
    lines: List[str] = []
    chunks: List[Dict] = []
    textless: List[int] = []
    pic_tmp_dirs: List[Path] = []

    for idx, slide in enumerate(prs.slides, start=1):
        page_lines = [f"## 第 {idx} 页"]
        has_text = False
        for shape in _iter_pptx_shapes(slide.shapes):
            st = shape.shape_type
            if st == MSO_SHAPE_TYPE.PICTURE:
                w_px = shape.width / 914400 * 96
                h_px = shape.height / 914400 * 96
                if w_px < _PPTX_IMAGE_MIN_PX or h_px < _PPTX_IMAGE_MIN_PX:
                    continue  # 装饰小图，丢弃
                png, tmpdir = _save_pptx_picture(shape, path)
                if png:
                    pic_tmp_dirs.append(tmpdir)
                    chunks.extend(_picture_chunks(png))
                    has_text = True
                continue
            if st == MSO_SHAPE_TYPE.CHART:
                md = _chart_to_markdown(shape.chart)
                if md:
                    chunks.append(
                        {
                            "text": md,
                            "metadata": json.dumps({"kind": "table", "page": idx}, ensure_ascii=False),
                        }
                    )
                has_text = True
                continue
            if st == MSO_SHAPE_TYPE.TABLE:
                md = _table_to_markdown(shape.table)
                if md:
                    chunks.append(
                        {
                            "text": md,
                            "metadata": json.dumps({"kind": "table", "page": idx}, ensure_ascii=False),
                        }
                    )
                has_text = True
                continue
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    page_lines.append(t)
                    has_text = True

        if has_text:
            lines.append("\n".join(page_lines))
        else:
            textless.append(idx)

    if textless:
        rendered = _render_pptx_pages(path, textless)
        for n, txt in rendered.items():
            lines.append(f"## 第 {n} 页\n{txt}")

    for d in pic_tmp_dirs:
        shutil.rmtree(d, ignore_errors=True)

    return "\n".join(lines), chunks


def _pdf_pages_ocr(path: Path) -> str:
    """扫描版 PDF 兜底：逐页渲染成 PNG，再用 RapidOCR 提取文字。

    用于 MinerU 失败且 PyMuPDF 提取不到内嵌文字（扫描件）的情况。
    失败返回空串，不阻断入库。
    """
    try:
        import fitz

        doc = fitz.open(str(path))
        tmp = _mkdtemp("pdf_ocr_", config.DATA_DIR)
        parts: List[str] = []
        try:
            for i, page in enumerate(doc, start=1):
                pix = page.get_pixmap(dpi=150)
                png = tmp / f"p{i}.png"
                pix.save(str(png))
                t = _ocr_image(png)
                if t:
                    parts.append(f"## 第 {i} 页\n{t}")
        finally:
            shutil.rmtree(tmp, ignore_errors=True)
        return "\n".join(parts)
    except Exception:  # noqa: BLE001 —— OCR 兜底失败不阻断
        return ""


def _extract_text(path: Path) -> str:
    ext = path.suffix.lower()
    if ext in {".md", ".markdown", ".txt"}:
        return path.read_text(encoding="utf-8", errors="ignore")
    if ext == ".docx":
        return _extract_docx(path)[0]
    if ext == ".pptx":
        return _extract_pptx(path)[0]
    if ext == ".pdf":
        try:
            return _pdf_to_markdown(path)
        except subprocess.CalledProcessError as e:
            logger.warning("MinerU 解析失败，回退 PyMuPDF 纯文本：%s", (e.stderr or "")[-2000:])
        except Exception as e:  # noqa: BLE001
            logger.warning("MinerU 解析失败，回退 PyMuPDF 纯文本：%s", e)
        import fitz  # PyMuPDF 兜底

        doc = fitz.open(str(path))
        text = "\n".join(page.get_text() for page in doc)
        if text.strip():
            return text
        # 扫描版 PDF：无内嵌文字 → 渲染每页 + RapidOCR 兜底
        ocr_text = _pdf_pages_ocr(path)
        if ocr_text.strip():
            logger.info("扫描版 PDF 已通过 RapidOCR 提取文字：%s", path.name)
            return ocr_text
        return text
    raise ValueError(f"不支持的文档类型: {ext}")


# ---------- 图片处理 ----------

def _describe_image(path: Path) -> str:
    """用 Qwen-VL 描述图片语义（DashScope API，无需本地模型）。"""
    b64 = base64.b64encode(path.read_bytes()).decode()
    msg = llm.generator.invoke(
        [
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": "详细描述这张图片的内容、关键文字和图表数据。"},
                    {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{b64}"}},
                ],
            }
        ]
    )
    return msg.content if isinstance(msg.content, str) else str(msg.content)


def _ocr_image(path: Path) -> str:
    """RapidOCR 抽取图片里的文字（首次调用会下载 onnx 模型）。失败返回空串。"""
    try:
        from rapidocr_onnxruntime import RapidOCR

        ocr = RapidOCR()
        result, _ = ocr(str(path))
        if not result:
            return ""
        return "\n".join(r[1] for r in result)
    except Exception:  # noqa: BLE001 —— OCR 失败不阻断入库
        return ""


def ingest_image(path: Path) -> int:
    """图片入库：OCR 文字 + VLM 描述分别作为一个 chunk。"""
    chunks = _picture_chunks(path)
    if not chunks:
        return 0
    return _embed_and_insert(chunks, path.name, "image")


# ---------- 核心入库 ----------

def _embed_and_insert(chunks: List[Dict], source: str, doc_type: str) -> int:
    if not chunks:
        return 0  # 空文档不入库，避免对空列表调 embedding 报错
    texts = [c["text"] for c in chunks]
    vectors = llm.embeddings.embed_documents(texts)
    data = [
        {
            "text": c["text"],
            "dense": v,
            "source": source,
            "doc_type": doc_type,
            "metadata": json.dumps(c["metadata"], ensure_ascii=False),
        }
        for c, v in zip(chunks, vectors)
    ]
    return store.insert(data)


def ingest_text(text: str, source: str = "inline", doc_type: str = "text") -> int:
    """纯文本入库。"""
    chunks = split_text(text)
    return _embed_and_insert(chunks, source, doc_type)


def ingest_file(path: str | Path) -> int:
    """按文件类型入库，返回写入的 chunk 数。"""
    path = Path(path)
    if not path.exists():
        raise FileNotFoundError(f"文件不存在: {path}")

    if path.suffix.lower() in {".png", ".jpg", ".jpeg", ".bmp", ".webp"}:
        return ingest_image(path)

    doc_type = path.suffix.lower().lstrip(".")

    if path.suffix.lower() == ".docx":
        md_text, table_chunks = _extract_docx(path)
        chunks = split_text(md_text) + table_chunks
    elif path.suffix.lower() == ".pptx":
        md_text, extra_chunks = _extract_pptx(path)
        chunks = split_text(md_text) + extra_chunks
    else:
        chunks = split_text(_extract_text(path))
    return _embed_and_insert(chunks, path.name, doc_type)
